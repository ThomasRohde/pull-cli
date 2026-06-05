from __future__ import annotations

import json
from pathlib import Path

TEXT_SUFFIXES = {".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".svg", ".log", ".yaml", ".yml"}


def extract_text_sidecar(path: Path) -> str | None:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return _read_text(path)
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    return None


def write_extracted_markdown(path: Path, text: str) -> Path:
    sidecar = path.with_name(f"{path.stem}.extracted.md")
    sidecar.write_text(f"# Extracted Text: {path.name}\n\n{text.strip()}\n", encoding="utf-8")
    return sidecar


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf(path: Path) -> str | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        return None
    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages).strip() or None


def _extract_docx(path: Path) -> str | None:
    try:
        import docx
    except ImportError:
        return None
    document = docx.Document(str(path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text).strip() or None


def _extract_xlsx(path: Path) -> str | None:
    try:
        import openpyxl
    except ImportError:
        return None
    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                lines.append("\t".join(values))
    return "\n".join(lines).strip() or None


def _extract_pptx(path: Path) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    presentation = Presentation(str(path))
    lines: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines.append(f"## Slide {slide_number}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text for cell in row.cells])
                lines.append(json.dumps(rows, ensure_ascii=False))
    return "\n".join(lines).strip() or None
